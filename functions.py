from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings(model='text-embedding-ada-002')
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI(model_name='gpt-3.5-turbo')
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(llm=llm, retriever=vectorstore.as_retriever(), memory=memory)
    return conversation_chain

def get_text_chunks_from_docx(docx_file):
    document = Document(docx_file)
    text = ""
    for para in document.paragraphs:
        text += para.text + "\n"
    return get_text_chunks(text)

def get_vectorstore_from_docx(text_chunks):
    embeddings = OpenAIEmbeddings(model='text-embedding-ada-002')
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain_for_docx(vectorstore):
    llm = ChatOpenAI(model_name='gpt-3.5-turbo')
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(llm=llm, retriever=vectorstore.as_retriever(), memory=memory)
    return conversation_chain

def get_text_chunks_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return get_text_chunks(text)

def get_vectorstore_from_ppt(text_chunks):
    embeddings = OpenAIEmbeddings(model='text-embedding-ada-002')
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain_for_ppt(vectorstore):
    llm = ChatOpenAI(model_name='gpt-3.5-turbo')
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(llm=llm, retriever=vectorstore.as_retriever(), memory=memory)
    return conversation_chain

def get_text_chunks_from_excel(excel_file):
    df = pd.read_excel(excel_file)
    text = ""
    for column in df.columns:
        text += f"{column}: {df[column].tolist()}\n"
    return get_text_chunks(text)

def get_vectorstore_from_excel(text_chunks):
    embeddings = OpenAIEmbeddings(model='text-embedding-ada-002')
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain_for_excel(vectorstore):
    llm = ChatOpenAI(model_name='gpt-3.5-turbo')
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(llm=llm, retriever=vectorstore.as_retriever(), memory=memory)
    return conversation_chain